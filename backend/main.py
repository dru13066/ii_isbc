import uuid
import json
from fastapi import FastAPI, UploadFile, File, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv

from pptx import Presentation

load_dotenv()

app = FastAPI(title="AI PPTX Builder")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "backend/uploads"

jobs = {}


@app.post("/generate")
async def generate(
    background_tasks: BackgroundTasks,
    template: UploadFile = File(...),
    content: UploadFile = File(...)
):
    job_id = str(uuid.uuid4())

    jobs[job_id] = {
        "status": "processing",
        "progress": 0,
        "file": None
    }

    template_path = f"{UPLOAD_DIR}/{job_id}_template.pptx"
    content_path = f"{UPLOAD_DIR}/{job_id}_content.json"
    output_path = f"{UPLOAD_DIR}/{job_id}_result.pptx"

    with open(template_path, "wb") as f:
        f.write(await template.read())

    with open(content_path, "wb") as f:
        f.write(await content.read())

    with open(content_path, "r", encoding="utf-8") as f:
        content_data = json.load(f)

    background_tasks.add_task(
        run_generation,
        job_id,
        template_path,
        content_data,
        output_path
    )

    return {"job_id": job_id}


def run_generation(job_id, template_path, content_data, output_path):
    total_slides = len(content_data["slides"])
    prs = Presentation(template_path)

    for index, slide_data in enumerate(content_data["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = slide_data["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

        jobs[job_id]["progress"] = int(((index + 1) / total_slides) * 100)

    prs.save(output_path)

    jobs[job_id]["status"] = "done"
    jobs[job_id]["file"] = output_path


@app.get("/status/{job_id}")
def status(job_id: str):
    job = jobs.get(job_id)
    if not job:
        return JSONResponse(status_code=404, content={"status": "not_found"})
    return job


@app.get("/download/{job_id}")
def download(job_id: str):
    job = jobs.get(job_id)
    if job and job["status"] == "done":
        return FileResponse(
            job["file"],
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="presentation.pptx"
        )
    return JSONResponse(status_code=404, content={"error": "File not ready"})
