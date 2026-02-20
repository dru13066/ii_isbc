from pptx import Presentation
from .ai import optimize_text

MAX_CHARS_PER_BULLET = 90

def generate_presentation(template_path, content, output_path):
    prs = Presentation(template_path)

    for slide_data in content["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = slide_data["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for bullet in slide_data["bullets"]:
            text = bullet
            if len(text) > MAX_CHARS_PER_BULLET:
                text = optimize_text(text, MAX_CHARS_PER_BULLET)

            p = tf.add_paragraph()
            p.text = text
            p.level = 0

    prs.save(output_path)